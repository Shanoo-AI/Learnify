from pptx import Presentation
import requests
import os
import re
import asyncio
import edge_tts  # Edge TTS

# ========== CONFIG ==========
PPT_PATH = "testing.pptx"      # <-- your hardcoded ppt file
OUTPUT_AUDIO = "explanation.mp3"

# OpenRouter / LLM API
API_KEY = os.getenv("OPENROUTER_API_KEY")
API_URL = "https://openrouter.ai/api/v1/chat/completions"
MODEL = "deepseek/deepseek-v3.2"

# Edge TTS voice
EDGE_VOICE = "ur-PK-UzmaNeural"  # Natural female voice. Change if you want male: "en-US-GuyNeural"

# ========== STEP 1: EXTRACT TEXT FROM PPT ==========
def extract_text_from_ppt(ppt_path):
    prs = Presentation(ppt_path)
    all_text = []

    for i, slide in enumerate(prs.slides, start=1):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    slide_text.append(text)

        if slide_text:
            all_text.append(f"Slide {i}:\n" + "\n".join(slide_text))

    return "\n\n".join(all_text)

# ========== STEP 2: SEND TO AI FOR EXPLANATION ==========
def generate_explanation(text):
    prompt = f"""
You are generating spoken lecture content for students.

Rules (must follow strictly):
Do NOT include any introduction like “Of course”, “Sure”, “I will help you”, “As an AI”, or “As a teacher”.
Do NOT mention yourself, the user, or the task.
Do NOT use bullet points, markdown, asterisks, headings, or special formatting.
Do NOT include meta commentary or instructions.
Write only plain, natural spoken English suitable for text-to-speech.
Start directly with the explanation content.
Keep the tone simple, clear, and student-friendly.

Now explain the following content: in urdu language

{text}
"""

    headers = {
        "Authorization": f"Bearer {API_KEY}",
        "Content-Type": "application/json",
    }

    payload = {
        "model": MODEL,
        "messages": [
            {"role": "user", "content": prompt}
        ],
        "temperature": 0.5
    }

    response = requests.post(API_URL, headers=headers, json=payload)
    response.raise_for_status()

    data = response.json()
    explanation = data["choices"][0]["message"]["content"]
    return explanation

# ========== CLEAN TEXT FOR TTS ==========
def clean_text_for_tts(text):
    text = re.sub(r"[*#_`>-]", " ", text)
    text = re.sub(r"-{2,}", " ", text)
    text = re.sub(r"\s+", " ", text)
    text = re.sub(r"[^\w\s.,!?]", " ", text)
    return text.strip()

# ========== STEP 3: TEXT TO SPEECH (EDGE TTS) ==========
async def text_to_audio_edge(text, output_file, voice=EDGE_VOICE):
    communicate = edge_tts.Communicate(text, voice)
    await communicate.save(output_file)

def text_to_audio(text, output_file):
    asyncio.run(text_to_audio_edge(text, output_file))

# ========== MAIN ==========
def main():
    print("📄 Reading PPT...")
    ppt_text = extract_text_from_ppt(PPT_PATH)

    if not ppt_text.strip():
        print("❌ No text found in PPT!")
        return

    print("🤖 Generating explanation from AI...")
    explanation = generate_explanation(ppt_text)

    clean_explanation = clean_text_for_tts(explanation)

    print("🔊 Converting explanation to audio with Edge TTS...")
    text_to_audio(clean_explanation, OUTPUT_AUDIO)

    print("✅ Done!")
    print(f"🎧 Audio saved as: {OUTPUT_AUDIO}")

if __name__ == "__main__":
    main()
