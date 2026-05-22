import asyncio
import os
import re
import shutil
import tempfile
from pathlib import Path

import edge_tts
import requests
from flask import after_this_request, jsonify, request, send_file, session
from pptx import Presentation
from werkzeug.utils import secure_filename

LANGUAGE_CONFIG = {
    "en": {
        "voice": "en-US-BrianMultilingualNeural",
        "instruction": "Explain in simple spoken English.",
    },
    "ur": {
        "voice": "ur-PK-UzmaNeural",
        "instruction": "Explain in simple spoken Urdu. Respond only in Urdu and do not use English.",
    },
}

ALLOWED_EXTENSIONS = {".ppt", ".pptx", ".pdf"}
OPENROUTER_URL = "https://openrouter.ai/api/v1/chat/completions"
DEFAULT_MODEL = "deepseek/deepseek-v3.2"
MAX_EXTRACTED_CHARS = 16000
OPENROUTER_API_KEY = os.getenv("OPENROUTER_API_KEY")


class AudioGenerationError(Exception):
    """Domain-level error for upload/LLM/TTS flow."""


def configure_audio_upload(app):
    app.config.setdefault("MAX_CONTENT_LENGTH", 30 * 1024 * 1024)


def register_audio_routes(app):
    configure_audio_upload(app)

    @app.route("/api/generate-audio", methods=["POST"])
    def generate_audio_from_upload():
        if not session.get("logged_in"):
            return jsonify({"success": False, "error": "Please login first"}), 401

        upload = request.files.get("file")
        if upload is None or not upload.filename:
            return jsonify({"success": False, "error": "No file uploaded"}), 400

        language = (request.form.get("language") or "en").lower()
        if language not in LANGUAGE_CONFIG:
            return jsonify({"success": False, "error": "Unsupported language"}), 400

        extension = Path(upload.filename).suffix.lower()
        if extension not in ALLOWED_EXTENSIONS:
            return jsonify(
                {"success": False, "error": "Only PPT, PPTX, and PDF files are allowed"}
            ), 400

        temp_dir = tempfile.mkdtemp(prefix="doc_audio_")
        input_path = os.path.join(temp_dir, secure_filename(upload.filename))
        output_path = os.path.join(temp_dir, "explanation.mp3")

        @after_this_request
        def cleanup(response):
            shutil.rmtree(temp_dir, ignore_errors=True)
            return response

        try:
            upload.save(input_path)
            if os.path.getsize(input_path) == 0:
                return jsonify({"success": False, "error": "Uploaded file is empty"}), 400

            extracted_text = extract_text_by_type(input_path, extension)
            if not extracted_text.strip():
                return jsonify({"success": False, "error": "No readable text found in file"}), 400

            explanation = generate_explanation_with_llm(extracted_text, language)
            cleaned_explanation = clean_text_for_tts(explanation)
            if not cleaned_explanation:
                return jsonify({"success": False, "error": "Generated explanation is empty"}), 502

            text_to_audio(cleaned_explanation, output_path, LANGUAGE_CONFIG[language]["voice"])
            return send_file(output_path, as_attachment=False, mimetype="audio/mpeg")

        except AudioGenerationError as exc:
            return jsonify({"success": False, "error": str(exc)}), 502
        except requests.RequestException:
            return jsonify({"success": False, "error": "LLM API request failed"}), 502
        except Exception:
            return jsonify({"success": False, "error": "Failed to process uploaded file"}), 500


def extract_text_by_type(file_path, extension):
    if extension in {".ppt", ".pptx"}:
        return extract_text_from_ppt(file_path)
    if extension == ".pdf":
        return extract_text_from_pdf(file_path)
    raise AudioGenerationError("Invalid file type")


def extract_text_from_ppt(ppt_path):
    presentation = Presentation(ppt_path)
    chunks = []

    for index, slide in enumerate(presentation.slides, start=1):
        slide_text = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "").strip()
            if text:
                slide_text.append(text)

        if slide_text:
            chunks.append(f"Slide {index}:\n" + "\n".join(slide_text))

    return "\n\n".join(chunks)[:MAX_EXTRACTED_CHARS]


def extract_text_from_pdf(pdf_path):
    extracted = extract_text_with_pymupdf(pdf_path)
    if extracted.strip():
        return extracted[:MAX_EXTRACTED_CHARS]

    extracted = extract_text_with_pdfplumber(pdf_path)
    return extracted[:MAX_EXTRACTED_CHARS]


def extract_text_with_pymupdf(pdf_path):
    try:
        import fitz
    except ImportError:
        return ""

    chunks = []
    with fitz.open(pdf_path) as doc:
        for index, page in enumerate(doc, start=1):
            page_text = page.get_text("text").strip()
            if page_text:
                chunks.append(f"Page {index}:\n{page_text}")
    return "\n\n".join(chunks)


def extract_text_with_pdfplumber(pdf_path):
    try:
        import pdfplumber
    except ImportError as exc:
        raise AudioGenerationError(
            "PDF support requires PyMuPDF or pdfplumber. Install one of them."
        ) from exc

    chunks = []
    with pdfplumber.open(pdf_path) as pdf:
        for index, page in enumerate(pdf.pages, start=1):
            page_text = (page.extract_text() or "").strip()
            if page_text:
                chunks.append(f"Page {index}:\n{page_text}")
    return "\n\n".join(chunks)


def build_prompt(content, language):
    instruction = LANGUAGE_CONFIG[language]["instruction"]
    return (
        "You are an expert teacher creating spoken lecture content for students.\n\n"

    "Strict rules (must follow):\n"
    "Do NOT read, repeat, or paraphrase the source text line by line.\n"
    "Do NOT sound like you are reading notes or slides aloud.\n"
    "Instead, TEACH the ideas as if explaining to a beginner in a classroom.\n"
    "Break down concepts into simpler terms, explain meanings, purposes, and examples where helpful.\n"
    "Connect ideas naturally like a teacher giving a lecture.\n"
    "If the source contains definitions, explain what they mean in practical or simple words.\n"
    "If the source contains lists or points, convert them into natural spoken explanation.\n\n"

    "Formatting rules:\n"
    "Do NOT include introductions like 'Sure', 'Of course', or 'Let's begin'.\n"
    "Do NOT mention yourself, the user, or the task.\n"
    "Do NOT use bullet points, markdown, headings, numbering, or special formatting.\n"
    "Write only plain natural spoken language suitable for text-to-speech.\n"
    "Start directly with the explanation.\n\n"

    f"Teaching style instruction:\n{instruction}\n\n"

    "Source material to explain:\n"
    f"{content}"
    )


def generate_explanation_with_llm(extracted_text, language):
    api_key = OPENROUTER_API_KEY
    if not api_key:
        raise AudioGenerationError("OPENROUTER_API_KEY is not configured")

    model = os.getenv("OPENROUTER_MODEL", DEFAULT_MODEL)
    headers = {
        "Authorization": f"Bearer {api_key}",
        "Content-Type": "application/json",
    }
    payload = {
        "model": model,
        "messages": [{"role": "user", "content": build_prompt(extracted_text, language)}],
        "temperature": 0.5,
    }

    response = requests.post(OPENROUTER_URL, headers=headers, json=payload, timeout=90)
    response.raise_for_status()
    data = response.json()

    try:
        return data["choices"][0]["message"]["content"].strip()
    except (KeyError, IndexError, TypeError) as exc:
        raise AudioGenerationError("Unexpected LLM API response format") from exc


def clean_text_for_tts(text):
    text = re.sub(r"[*#_`>-]+", " ", text)
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{2,}", "\n", text)
    # Preserve Unicode letters (including Urdu) and common punctuation.
    text = re.sub(r"[^\w\s.,!?،۔؟:;'\"]", " ", text, flags=re.UNICODE)
    return text.strip()


async def text_to_audio_edge(text, output_file, voice):
    communicator = edge_tts.Communicate(text, voice)
    await communicator.save(output_file)


def text_to_audio(text, output_file, voice):
    asyncio.run(text_to_audio_edge(text, output_file, voice))
