from pptx import Presentation
from gtts import gTTS
import requests
import os

# ========== CONFIG ==========
PPT_PATH = "testing.pptx"      # <-- your hardcoded ppt file
OUTPUT_AUDIO = "explanation.mp3"

# Set OPENROUTER_API_KEY in your environment.
API_KEY = os.getenv("OPENROUTER_API_KEY")
API_URL = "https://openrouter.ai/api/v1/chat/completions"
MODEL = "deepseek/deepseek-v3.2"   # or any model you want

# ========== STEP 1: EXTRACT TEXT FROM PPT ==========
def extract_text_from_ppt(ppt_path):
    prs = Presentation(ppt_path)
    all_text = []

    for i, slide in enumerate(prs.slides, start=1):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    slide_text.append(text)

        if slide_text:
            all_text.append(f"Slide {i}:\n" + "\n".join(slide_text))

    return "\n\n".join(all_text)

# ========== STEP 2: SEND TO AI FOR EXPLANATION ==========
def generate_explanation(text):
    prompt = f"""
    You are a helpful teacher. Explain the following presentation content
in simple, clear, student-friendly language, slide by slide
You are generating spoken lecture content for students.

Your task: Explain the following content clearly and simply, as if a teacher is explaining it in class.

Rules (must follow strictly):

Do NOT include any introduction like “Of course”, “Sure”, “I will help you”, “As an AI”, or “As a teacher”.

Do NOT mention yourself, the user, or the task.

Do NOT use bullet points, markdown, asterisks, headings, or special formatting.

Do NOT include meta commentary or instructions.

Write only plain, natural spoken English suitable for text-to-speech.

Start directly with the explanation content.

Keep the tone simple, clear, and student-friendly.

Now explain the following content:

{text}
"""

    headers = {
        "Authorization": f"Bearer {API_KEY}",
        "Content-Type": "application/json",
    }

    payload = {
        "model": MODEL,
        "messages": [
            {"role": "user", "content": prompt}
        ],
        "temperature": 0.5
    }

    response = requests.post(API_URL, headers=headers, json=payload)
    response.raise_for_status()

    data = response.json()
    explanation = data["choices"][0]["message"]["content"]
    return explanation

# ========== STEP 3: TEXT TO SPEECH ==========
def text_to_audio(text, output_file):
    tts = gTTS(text=text, lang="en")
    tts.save(output_file)

# ========== MAIN ==========
def main():
    print("📄 Reading PPT...")
    ppt_text = extract_text_from_ppt(PPT_PATH)

    if not ppt_text.strip():
        print("❌ No text found in PPT!")
        return

    print("🤖 Generating explanation from AI...")
    explanation = generate_explanation(ppt_text)

    print("🔊 Converting explanation to audio...")
    text_to_audio(explanation, OUTPUT_AUDIO)

    print("✅ Done!")
    print(f"🎧 Audio saved as: {OUTPUT_AUDIO}")

if __name__ == "__main__":
    main()
